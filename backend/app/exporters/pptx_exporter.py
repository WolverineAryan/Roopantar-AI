from pathlib import Path
from typing import Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def export_presentation_to_pptx(content: Dict[str, Any], output_path: Path) -> Path:
    prs = Presentation()
    
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slides_data = content.get("slides", [])
    deck_title = content.get("deck_title", "Roopantar-AI Executive Presentation")
    subtitle = content.get("subtitle", "Strategic Analysis & Multi-Format Intelligence")
    
    # Modern Brand Palette (Slate 900, Electric Pink, Radiant Coral)
    primary_color = RGBColor(15, 23, 42)      # Slate 900
    card_bg_color = RGBColor(248, 250, 252)   # Slate 50
    accent_pink = RGBColor(236, 72, 153)      # Hot Pink (#EC4899)
    accent_purple = RGBColor(124, 58, 237)    # Purple (#7C3AED)
    text_dark = RGBColor(15, 23, 42)          # Slate 900 for high-contrast readability
    text_muted = RGBColor(71, 85, 105)        # Slate 600
    white_color = RGBColor(255, 255, 255)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_slide_layout)
        
        slide_title = slide_data.get("title", f"Slide {i+1}")
        bullet_points = slide_data.get("bullet_points", [])
        speaker_notes = slide_data.get("speaker_notes", "")
        slide_type = slide_data.get("slide_type", "Content")
        
        # Add Speaker Notes to PPTX Notes Frame
        if speaker_notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = f"Presenter Script:\n{speaker_notes}"
            
        if i == 0 or slide_type == "Title":
            # Title Slide Layout - Premium Dark Executive Canvas
            bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = primary_color
            bg_shape.line.fill.background()
            
            # Left Glowing Brand Accent Bar
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.8), Inches(0.18), Inches(4.0))
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent_pink
            bar.line.fill.background()
            
            # Title & Subtitle Text Box
            txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.6), Inches(10.8), Inches(4.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            
            # Slide Title (Pt 48 - Large, bold, commanding)
            p = tf.paragraphs[0]
            p.text = slide_title
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = white_color
            p.space_after = Pt(12)
            
            # Subtitle (Pt 24)
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(24)
            p2.font.color.rgb = RGBColor(226, 232, 240) # Slate 200
            p2.space_after = Pt(20)
            
            # Bullet Highlights if present
            if bullet_points:
                for bp in bullet_points[:3]:
                    p3 = tf.add_paragraph()
                    p3.text = f"•  {bp}"
                    p3.font.size = Pt(18)
                    p3.font.color.rgb = RGBColor(203, 213, 225)
                    p3.space_after = Pt(8)
                    
            # Footer on Title Slide
            footer_box = slide.shapes.add_textbox(Inches(1.5), Inches(6.5), Inches(10.5), Inches(0.5))
            tf_footer = footer_box.text_frame
            p_foot = tf_footer.paragraphs[0]
            p_foot.text = "Roopantar-AI • Enterprise Intelligence Platform"
            p_foot.font.size = Pt(12)
            p_foot.font.color.rgb = RGBColor(148, 163, 184)
                
        else:
            # Content Slide Layout - Bright Modern Executive Design
            
            # Top Banner Background
            top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.6))
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = primary_color
            top_bar.line.fill.background()
            
            # Pink Accent Strip under top banner
            strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.55), Inches(13.333), Inches(0.06))
            strip.fill.solid()
            strip.fill.fore_color.rgb = accent_pink
            strip.line.fill.background()
            
            # Slide Header Tag & Title Text Box
            title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.25), Inches(11.333), Inches(1.2))
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            
            p_tag = tf_title.paragraphs[0]
            p_tag.text = (slide_type or "ANALYSIS & DIRECTIVES").upper()
            p_tag.font.size = Pt(11)
            p_tag.font.bold = True
            p_tag.font.color.rgb = accent_pink
            p_tag.space_after = Pt(4)
            
            p_title = tf_title.add_paragraph()
            p_title.text = slide_title
            p_title.font.size = Pt(32)
            p_title.font.bold = True
            p_title.font.color.rgb = white_color
            
            # Main Content Container Box
            content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.6))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            
            # Large, Readable, High-Contrast Bullet Points (Pt 24 with 1.35 Line Spacing)
            for j, bp in enumerate(bullet_points):
                p = tf_content.paragraphs[0] if j == 0 else tf_content.add_paragraph()
                p.text = f"•   {bp}"
                p.font.size = Pt(24)
                p.font.bold = False
                p.font.color.rgb = text_dark
                p.space_after = Pt(22)
                p.line_spacing = 1.35
                
            # Footer Details
            footer_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.85), Inches(11.333), Inches(0.4))
            tf_footer = footer_box.text_frame
            p_footer = tf_footer.paragraphs[0]
            p_footer.text = f"Roopantar-AI • Executive Briefing  |  Slide {i+1} of {len(slides_data)}"
            p_footer.font.size = Pt(11)
            p_footer.font.color.rgb = text_muted

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
